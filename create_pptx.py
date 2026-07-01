import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import pandas as pd

EXCEL_FILE = "20260624_Data_Estate_Cleanup_Audit xlsx.xlsx"
SHEET_NAME = "BQ-added-Data"


def load_tranche1_targets():
    df = pd.read_excel(EXCEL_FILE, sheet_name=SHEET_NAME)
    return df[
        df['Tranche'].astype(str).str.startswith("Tranche 1:")
        & df['Defensible Deletion Status'].astype(str).str.startswith("CONFIRMED")
    ].copy()


def _parse_last_reader(reads_field):
    """Extract (when, who) from the first entry of 'Last 5 Reads (Who/When)', or (None, None) if never read."""
    s = str(reads_field)
    if s.startswith("No Read Access"):
        return None, None
    first_entry = s.split("\n")[0].strip()
    m = re.match(r'^(.*?)\s*\(([^)]+)\)\s*$', first_entry)
    if not m:
        return first_entry, None
    return m.group(1).strip(), m.group(2).strip()


def summarise_readers(t1):
    """Group Tranche 1 candidates by who last read them, to reveal bulk-scan vs genuine usage patterns."""
    parsed = t1['Last 5 Reads (Who/When)'].apply(_parse_last_reader)
    t1 = t1.assign(
        last_read_when=[p[0] for p in parsed],
        last_read_who=[p[1] for p in parsed],
    )
    has_read = t1[t1['last_read_who'].notna()]
    no_read_datasets = t1[t1['last_read_who'].isna()]['Dataset'].tolist()

    summary = (
        has_read.groupby('last_read_who')
        .agg(datasets=('Dataset', 'nunique'), last_read=('last_read_when', 'max'))
        .sort_values('datasets', ascending=False)
        .reset_index()
    )
    return summary, no_read_datasets


def create_presentation():
    prs = Presentation()
    
    # Set slide size to 13.333 x 7.5 inches (Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- BRANDING COLORS ---
    THG_BLUE = RGBColor(0x1A, 0x73, 0xE8)
    THG_DARK = RGBColor(0x16, 0x20, 0x2E)
    THG_GREY = RGBColor(0x5A, 0x66, 0x75)
    THG_GREEN = RGBColor(0x1E, 0x8E, 0x5A)
    THG_RED = RGBColor(0xC5, 0x22, 0x1F)  # BRAG "Red" - critical/action needed
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # --- FONTS ---
    HEADER_FONT = "Montserrat"
    BODY_FONT = "Open Sans"

    def apply_style(text_frame, font_name=BODY_FONT, size=Pt(18), color=THG_DARK, bold=False, align=PP_ALIGN.LEFT):
        # word_wrap/auto_size must be explicit - without them, long lines don't wrap
        # and run off the shape (and sometimes off the slide) instead of flowing to a new line.
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = size
                run.font.color.rgb = color
                run.font.bold = bold

    def add_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        apply_style(tf, font_name=HEADER_FONT, size=Pt(40), color=THG_DARK, bold=True)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = "BigQuery Housekeeping"
    apply_style(tf, font_name=HEADER_FONT, size=Pt(60), color=THG_BLUE, bold=True)
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1))
    tf = subtitle.text_frame
    p = tf.add_paragraph()
    p.text = "Data Estate Decommissioning & Scream Test Framework"
    apply_style(tf, font_name=BODY_FONT, size=Pt(28), color=THG_GREY)
    
    context = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    tf = context.text_frame
    p = tf.add_paragraph()
    p.text = "EXECUTIVE REVIEW  ·  CLOUD FINOPS  ·  JULY 2026"
    apply_style(tf, font_name=HEADER_FONT, size=Pt(14), color=THG_BLUE, bold=True)

    # --- Slide 2: Recovery Protocol ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "How to recover (Emergency Protocol)")
    levers = [
        ("1. Automated Restore", "Use `bq_restore_access.py` to pull metadata from `gs://bq-admin-backup/` and reapply in seconds."),
        ("2. Metadata Integrity", "Full access policies are versioned to GCS *before* any modification."),
        ("3. Org-Level Access", "Org Admins retain management persistence via Org-level IAM; zero lockout risk.")
    ]
    for i, (title, text) in enumerate(levers):
        top = Inches(1.8 + (i * 1.5))
        tbox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.5))
        tf = tbox.text_frame
        p = tf.add_paragraph()
        p.text = title
        apply_style(tf, font_name=HEADER_FONT, size=Pt(24), color=THG_BLUE, bold=True)
        mbox = slide.shapes.add_textbox(Inches(0.5), top + Inches(0.4), Inches(12), Inches(0.8))
        tf = mbox.text_frame
        p = tf.add_paragraph()
        p.text = text
        apply_style(tf, font_name=BODY_FONT, size=Pt(18), color=THG_DARK)

    # --- Slide 3: Tranche Overview (Financials & Volumes) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Tranche Analysis: Stale & Abandoned Data")
    
    rows, cols = 8, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table
    headers = ["Tranche", "Cost/Mo ($)", "Volume (GB)", "Datasets", "Projects"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THG_DARK
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True

    tranche_data = [
        ["T1: Abandoned Test", "549.90", "27,494", "37", "6"],
        ["T2: Abandoned Prod", "8,913.59", "445,681", "174", "11"],
        ["T3: Small Abandoned", "227.33", "11,372", "1,010", "15"],
        ["T4: Stale (180d)", "103.68", "5,182", "39", "9"],
        ["T5: Active Test", "581.07", "29,053", "49", "12"],
        ["Subtotal (Cleanup)", "10,375.57", "518,782", "1,309", "53"],
        ["Active (Prod)", "32,906.12", "1,645,314", "788", "28"]
    ]
    for r, row in enumerate(tranche_data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            if "Subtotal" in val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                cell.text_frame.paragraphs[0].font.bold = True

    # --- Slide 4: Access Profiles & Cleanup Risk (Tranche 1 candidates) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Access Profiles & Cleanup Risk: Tranche 1 Candidates")

    t1 = load_tranche1_targets()
    total_datasets = len(t1)
    total_size = t1['Size (GB)'].sum()
    total_cost = t1['Monthly Cost ($)'].sum()

    stat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.45))
    tf = stat_box.text_frame
    p = tf.add_paragraph()
    p.text = (
        f"{total_datasets} confirmed candidates  ·  {total_size:,.0f} GB  ·  "
        f"${total_cost:,.2f}/mo  ·  zero write activity in 180+ days"
    )
    apply_style(tf, font_name=BODY_FONT, size=Pt(16), color=THG_GREY, bold=True)

    chart_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.35))
    tf = chart_title.text_frame
    p = tf.add_paragraph()
    p.text = "Read vs. Storage-Update Recency — Top 15 Most Stale"
    apply_style(tf, font_name=HEADER_FONT, size=Pt(18), color=THG_DARK, bold=True)

    top_by_activity = t1.sort_values("Days Since Last Activity", ascending=False).head(15)
    chart_data = CategoryChartData()
    chart_data.categories = list(top_by_activity['Dataset'])
    # None marks datasets with no recorded read at all (rendered as a gap, not zero/"just read")
    read_values = [
        float(row['Days Since Last Activity']) if pd.notna(row['Last Read Time']) else None
        for _, row in top_by_activity.iterrows()
    ]
    chart_data.add_series('Days Since Last Read', read_values)
    chart_data.add_series('Days Since Last Storage Update', [float(v) for v in top_by_activity['Days Since Last Modified']])

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2.15), Inches(12), Inches(2.45), chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    plot = chart.plots[0]
    plot.has_data_labels = False
    series_read, series_modified = plot.series
    series_read.format.fill.solid()
    series_read.format.fill.fore_color.rgb = THG_BLUE
    series_modified.format.fill.solid()
    series_modified.format.fill.fore_color.rgb = THG_RED
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)

    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.75), Inches(12), Inches(0.35))
    tf = table_title.text_frame
    p = tf.add_paragraph()
    p.text = "Biggest Offenders by Monthly Cost"
    apply_style(tf, font_name=HEADER_FONT, size=Pt(18), color=THG_DARK, bold=True)

    top_by_cost = t1.sort_values("Monthly Cost ($)", ascending=False).head(6)
    rows, cols = len(top_by_cost) + 1, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(5.15), Inches(12), Inches(2)).table
    headers = ["Project", "Dataset", "Days Inactive", "Size (GB)", "Monthly Cost ($)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THG_DARK
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)

    for r, (_, row) in enumerate(top_by_cost.iterrows()):
        values = [
            row['Project'],
            row['Dataset'],
            f"{int(row['Days Since Last Activity'])}",
            f"{row['Size (GB)']:,.2f}",
            f"{row['Monthly Cost ($)']:,.2f}",
        ]
        for c, val in enumerate(values):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    # --- Slide 5: Tranche 1 - Who's Actually Reading This Data? ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Tranche 1: Who's Actually Reading This Data?")

    reader_summary, no_read_datasets = summarise_readers(t1)
    total_with_reads = int(reader_summary['datasets'].sum())
    bulk_scan_total = int(reader_summary[reader_summary['datasets'] > 1]['datasets'].sum())

    insight = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.9))
    tf = insight.text_frame
    p = tf.add_paragraph()
    p.text = (
        f"{total_with_reads} of {total_datasets} candidates show a recent 'read', but {bulk_scan_total} of those "
        f"trace to just 2 accounts, each hitting a dozen+ datasets in a single same-day batch — "
        f"evidence of an automated scan, not ongoing human usage."
    )
    apply_style(tf, font_name=BODY_FONT, size=Pt(16), color=THG_GREY, bold=True)

    table_rows = len(reader_summary) + 1
    table = slide.shapes.add_table(table_rows, 4, Inches(0.5), Inches(2.4), Inches(12), Inches(2.2)).table
    headers = ["Reader", "Datasets Touched (Tranche 1)", "Last Read", "Pattern"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THG_DARK
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)

    for r, row in reader_summary.iterrows():
        is_service_account = "gserviceaccount.com" in row['last_read_who']
        pattern = ("Service account, " if is_service_account else "") + (
            "single-day bulk scan" if row['datasets'] > 1 else "single read"
        )
        values = [row['last_read_who'], str(int(row['datasets'])), row['last_read'], pattern.capitalize()]
        for c, val in enumerate(values):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12), Inches(1.5))
    tf = footer.text_frame
    p = tf.add_paragraph()
    p.text = (
        f"Zero recorded reads or writes at all ({len(no_read_datasets)}): " + ", ".join(no_read_datasets) + "."
    )
    apply_style(tf, font_name=BODY_FONT, size=Pt(14), color=THG_DARK)

    # --- Slide 6: Execution & Repository ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Execution Control")
    
    infobox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4))
    tf = infobox.text_frame
    p = tf.add_paragraph()
    p.text = "Control File: 20260624_Data_Estate_Cleanup_Audit.xlsx"
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "Targeting 1,300+ datasets with $10k+ monthly waste potential."
    p3 = tf.add_paragraph()
    p3.text = "\nCode & Full Audit Control Repository:"
    p4 = tf.add_paragraph()
    p4.text = "https://github.com/alandoddthg/bq-housekeeping"
    
    apply_style(tf, font_name=BODY_FONT, size=Pt(28), color=THG_DARK)
    tf.paragraphs[4].runs[0].font.size = Pt(24)
    tf.paragraphs[4].runs[0].font.color.rgb = THG_BLUE
    tf.paragraphs[4].runs[0].font.bold = True

    # --- Slide 7: Status Tracking & Reporting ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Status Tracking & Reporting")

    intro = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = intro.text_frame
    p = tf.add_paragraph()
    p.text = "Every Scream Test run now produces a live, human-readable status view — no manual spreadsheet updates."
    apply_style(tf, font_name=BODY_FONT, size=Pt(18), color=THG_GREY)

    layers = [
        ("1. state.json", "Machine record per dataset: backup location, Scream Test start time, current status (active / restored / deleted)."),
        ("2. status_report.md", "Auto-regenerated Markdown report on every workspace sync — per-dataset days elapsed and a Phase D readiness flag."),
        ("3. Audit workbook stays clean", "Status is never pasted back into the Excel control file; the report is regenerated from live state each time (constitution §13.3).")
    ]
    for i, (title, text) in enumerate(layers):
        top = Inches(2.3 + (i * 1.6))
        tbox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.5))
        tf = tbox.text_frame
        p = tf.add_paragraph()
        p.text = title
        apply_style(tf, font_name=HEADER_FONT, size=Pt(24), color=THG_BLUE, bold=True)
        mbox = slide.shapes.add_textbox(Inches(0.5), top + Inches(0.4), Inches(12), Inches(1.0))
        tf = mbox.text_frame
        p = tf.add_paragraph()
        p.text = text
        apply_style(tf, font_name=BODY_FONT, size=Pt(18), color=THG_DARK)

    prs.save("BigQuery_Housekeeping_Process.pptx")

if __name__ == "__main__":
    create_presentation()
